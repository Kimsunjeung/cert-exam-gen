import os
from pathlib import Path
from fastapi import UploadFile
import PyPDF2
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import pytesseract

class FileHandler:
    def __init__(self, downloads_dir: Path):
        self.downloads_dir = downloads_dir
    
    async def save_file(self, file: UploadFile) -> Path:
        """파일을 downloads 디렉터리에 저장"""
        file_path = self.downloads_dir / file.filename
        
        with open(file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
        
        print(f"✅ 파일 저장됨: {file_path}")
        return file_path
    
    async def extract_text(self, file_path: Path) -> str:
        """파일 형식에 따라 텍스트 추출"""
        extension = file_path.suffix.lower()
        
        try:
            if extension == ".pdf":
                return self._extract_from_pdf(file_path)
            elif extension in [".doc", ".docx"]:
                return self._extract_from_docx(file_path)
            elif extension in [".xls", ".xlsx"]:
                return self._extract_from_excel(file_path)
            elif extension in [".ppt", ".pptx"]:
                return self._extract_from_pptx(file_path)
            elif extension in [".jpg", ".jpeg", ".png"]:
                return self._extract_from_image(file_path)
            else:
                raise ValueError(f"지원하지 않는 파일 형식: {extension}")
        except Exception as e:
            print(f"텍스트 추출 오류: {e}")
            raise
    
    def _extract_from_pdf(self, file_path: Path) -> str:
        """PDF에서 텍스트 추출"""
        text = ""
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
    
    def _extract_from_docx(self, file_path: Path) -> str:
        """Word 문서에서 텍스트 추출"""
        doc = Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    
    def _extract_from_excel(self, file_path: Path) -> str:
        """Excel에서 텍스트 추출"""
        workbook = load_workbook(file_path)
        text = ""
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text += " ".join([str(cell) for cell in row if cell]) + "\n"
        return text
    
    def _extract_from_pptx(self, file_path: Path) -> str:
        """PowerPoint에서 텍스트 추출"""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def _extract_from_image(self, file_path: Path) -> str:
        """이미지에서 OCR로 텍스트 추출"""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang='kor+eng')
            return text
        except Exception as e:
            print(f"OCR 오류: {e}")
            return "이미지에서 텍스트를 추출할 수 없습니다."
